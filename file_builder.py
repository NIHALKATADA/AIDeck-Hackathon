import pandas as pd
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import io

def hex_to_rgb(hex_code):
    hex_code = hex_code.lstrip('#')
    r, g, b = tuple(int(hex_code[i:i+2], 16) for i in (0, 2, 4))
    return RGBColor(r, g, b)

def create_pptx(presentation_data, topic_name, theme_name, font_name, title_color_hex, bg_color_hex, text_color_hex):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    theme_title_color = hex_to_rgb(title_color_hex)
    theme_bg_color = hex_to_rgb(bg_color_hex)
    
    if theme_name == "Dark Executive":
        theme_text_color = RGBColor(255, 255, 255)
    else:
        theme_text_color = hex_to_rgb(text_color_hex)

    def apply_slide_background(slide, color_rgb):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color_rgb

    # --- TITLE SLIDE ---
    title_slide_layout = prs.slide_layouts[6] 
    slide = prs.slides.add_slide(title_slide_layout)
    apply_slide_background(slide, theme_bg_color)
    
    if theme_name == "Breeze":
        c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(3.5), Inches(6), Inches(6))
        c1.fill.solid()
        c1.fill.fore_color.rgb = hex_to_rgb("#E0F2FE")
        c1.line.fill.background()
    elif theme_name == "Dark Executive":
        c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-3), Inches(-1), Inches(8), Inches(8))
        c1.fill.solid()
        c1.fill.fore_color.rgb = hex_to_rgb("#1E3A8A")
        c1.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(2))
    p = title_box.text_frame.paragraphs[0]
    p.text = topic_name.title()
    p.font.name, p.font.size, p.font.bold = font_name, Pt(54), True
    p.font.color.rgb = theme_title_color

    # --- CONTENT SLIDES ---
    for slide_data in presentation_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_slide_background(slide, theme_bg_color)
        
        t_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(1))
        t_p = t_box.text_frame.paragraphs[0]
        t_p.text = slide_data.title
        t_p.font.name, t_p.font.size, t_p.font.bold = font_name, Pt(36), True
        t_p.font.color.rgb = theme_title_color

        start_y = Inches(1.6)
        for i, bullet in enumerate(slide_data.bullets[:6]):
            col, row = i % 2, i // 2
            x_pos = Inches(0.8) + (col * Inches(6.2))
            y_pos = start_y + (row * Inches(1.9))
            
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, y_pos, Inches(5.8), Inches(1.6))
            card.fill.solid()
            card.fill.fore_color.rgb = hex_to_rgb("#FFFFFF") if theme_name == "Breeze" else hex_to_rgb("#1E293B")
            card.line.fill.background()
            
            tb = slide.shapes.add_textbox(x_pos + Inches(0.2), y_pos + Inches(0.1), Inches(5.4), Inches(1.4))
            tb.text_frame.word_wrap = True
            p = tb.text_frame.paragraphs[0]
            p.text = bullet
            p.font.name, p.font.size, p.font.color.rgb = font_name, Pt(14), theme_text_color
            
    out_io = io.BytesIO()
    prs.save(out_io)
    out_io.seek(0)
    return out_io

def create_xlsx(dataset_data):
    df = pd.DataFrame([row.model_dump() for row in dataset_data])
    xlsx_io = io.BytesIO()
    with pd.ExcelWriter(xlsx_io, engine='openpyxl') as writer:
        df.to_excel(writer, index=False, sheet_name='AI Dataset')
    xlsx_io.seek(0)
    return xlsx_io